"""Build the canonical proposal template as a .pptx file.

Generates a 17-slide deck following the structure documented at
docs/next-60-days/proposal-template.md. The deck is profile-agnostic - placeholder
text in [SQUARE BRACKETS] is meant to be replaced by the FC for each prospect.

Run: python3 .tmp/build-proposal-template.py
Output: .tmp/proposal-template.pptx (then upload to Drive)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).resolve().parent / "proposal-template.pptx"

# AIA-flavoured brand palette (subdued — slides should not look gimmicky)
RED = RGBColor(0xC8, 0x10, 0x2E)         # AIA red
INK = RGBColor(0x1A, 0x1A, 0x1A)         # near-black text
MUTED = RGBColor(0x66, 0x66, 0x66)       # secondary text
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)       # card background
ACCENT = RGBColor(0xE8, 0xC8, 0x4E)      # gold accent for callouts
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_box(slide, x, y, w, h, fill=LIGHT, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_accent_strip(slide, x, y, w, color=RED, h=0.06):
    add_box(slide, x, y, w, h, fill=color)


def add_footer(slide, page_num, total=17):
    add_text(slide, 0.5, 7.1, 7, 0.3, "Proposal Template — built from the team's MDRT shape", size=9, color=MUTED)
    add_text(slide, 11.5, 7.1, 1.3, 0.3, f"{page_num} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------- Slide 1: Cover
s = prs.slides.add_slide(BLANK)
add_box(s, 0, 0, 13.333, 7.5, fill=WHITE)
add_accent_strip(s, 0.5, 1.0, 1.5, RED, h=0.08)
add_text(s, 0.5, 1.3, 12, 1.0, "[FINANCIAL REVIEW / RETIREMENT PLANNING / PROTECTION REVIEW]", size=14, color=MUTED)
add_text(s, 0.5, 2.2, 12, 1.4, "[Title that names the OUTCOME, not the product]", size=42, bold=True, color=INK)
add_text(s, 0.5, 3.7, 12, 0.5, "Specially For:", size=18, color=MUTED)
add_text(s, 0.5, 4.2, 12, 0.8, "[Client first name]", size=36, bold=True, color=RED)
add_text(s, 0.5, 6.5, 12, 0.4, "[Date] · Prepared by [Your name] · AIA Singapore", size=12, color=MUTED)
add_footer(s, 1)


# ---------------------------------------------------------------- Slide 2: Agenda
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "Agenda", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "What we'll cover today, in roughly 30 minutes.", size=14, color=MUTED)

agenda = [
    ("1.", "Where you are today", "Existing portfolio + life-stage gap analysis"),
    ("2.", "What continuing on the current path costs", "The cost of doing nothing, quantified"),
    ("3.", "Two options that improve the outcome", "Side-by-side, you choose"),
    ("4.", "The numbers behind both", "Track record, projections, timing"),
    ("5.", "Next steps if any of this resonates", "No pressure — your call"),
]
for i, (n, h, sub) in enumerate(agenda):
    y = 2.2 + i * 0.85
    add_text(s, 0.6, y, 0.5, 0.6, n, size=22, bold=True, color=RED)
    add_text(s, 1.2, y, 11, 0.4, h, size=18, bold=True, color=INK)
    add_text(s, 1.2, y + 0.42, 11, 0.4, sub, size=13, color=MUTED)
add_footer(s, 2)


# ---------------------------------------------------------------- Slide 3: The Framework
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "[Outcome] Planning Framework", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "A four-phase approach. Diagnose first; optimise second.", size=14, color=MUTED)

# Phase grouping labels
add_text(s, 1.2, 2.05, 5, 0.4, "Phase 1 — Diagnose", size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s, 7.0, 2.05, 5, 0.4, "Phase 2 — Optimise", size=12, bold=True, color=RED, align=PP_ALIGN.CENTER)

phases = [
    (1, "Readiness Assessment", "[Coverage adequacy / retirement readiness / asset allocation review]"),
    (2, "Optimisation Targets", "[Identify redundant / underperforming components]"),
    (3, "Tailored Strategy", "[Recommendations balancing risk + return]"),
    (4, "Execution + Review", "Execute if comfortable; review yearly"),
]
card_w = 2.85
gap = 0.15
start_x = 0.7
for i, (n, h, sub) in enumerate(phases):
    x = start_x + i * (card_w + gap)
    add_box(s, x, 2.5, card_w, 3.3, fill=LIGHT)
    add_box(s, x, 2.5, card_w, 0.5, fill=RED)
    add_text(s, x, 2.55, card_w, 0.4, str(n), size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.15, 3.15, card_w - 0.3, 0.6, h, size=15, bold=True, color=INK)
    add_text(s, x + 0.15, 3.85, card_w - 0.3, 1.7, sub, size=11, color=MUTED)
add_footer(s, 3)


# ---------------------------------------------------------------- Slide 4: Fact-find recap (FORM)
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "Where you are — the picture I heard", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "From our conversation. Correct me if I got anything wrong.", size=14, color=MUTED)

form = [
    ("F", "Family", "[2 kids (5, 8); spouse works part-time; parents in their 60s]"),
    ("O", "Occupation", "[SaaS sales, $180K/year, target $250K by 35]"),
    ("R", "Recreation", "[Triathlon, family travel, hiking]"),
    ("M", "Money", "[$80K savings, no investment plan, $400K mortgage, 1 ILP]"),
]
for i, (letter, label, content) in enumerate(form):
    row = i // 2
    col = i % 2
    x = 0.7 + col * 6.1
    y = 2.2 + row * 2.4
    add_box(s, x, y, 5.8, 2.1, fill=LIGHT)
    add_box(s, x, y, 0.7, 2.1, fill=RED)
    add_text(s, x, y + 0.55, 0.7, 1.0, letter, size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.85, y + 0.2, 4.8, 0.5, label, size=18, bold=True, color=INK)
    add_text(s, x + 0.85, y + 0.75, 4.8, 1.2, content, size=12, color=MUTED)
add_footer(s, 4)


# ---------------------------------------------------------------- Slide 5: Existing portfolio summary
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "Your existing portfolio — consolidated view", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "Most clients have never seen this on one page. This is your reference document.", size=14, color=MUTED)

# Header row
headers = ["Plan", "Type", "Sum Assured", "Premium", "Riders / Notes"]
widths = [3.2, 2.2, 1.8, 1.5, 3.5]
x = 0.5
for i, h in enumerate(headers):
    add_box(s, x, 2.1, widths[i], 0.5, fill=INK)
    add_text(s, x + 0.1, 2.15, widths[i] - 0.2, 0.4, h, size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    x += widths[i]

rows = [
    ["[AIA Pro Lifetime]", "[ILP whole-life]", "[$300K]", "[$600/mo]", "[CI rider; $42K cash value]"],
    ["[AIA Solitaire PA]", "[Personal Accident]", "[$750K]", "[$33/mo]", "[Plan 3]"],
    ["[HSGM Standard]", "[Hospital]", "[Standard]", "[Medisave]", "[No rider — $5K co-pay exposure]"],
    ["[CPF Life]", "[Annuity]", "[—]", "[—]", "[FRS, payout from 65]"],
    ["TOTALS", "", "[$1.05M]", "[$633/mo]", ""],
]
for r_idx, row in enumerate(rows):
    y = 2.6 + r_idx * 0.55
    is_total = r_idx == len(rows) - 1
    x = 0.5
    for c_idx, cell in enumerate(row):
        add_box(s, x, y, widths[c_idx], 0.55, fill=ACCENT if is_total else LIGHT if r_idx % 2 == 0 else WHITE, line=MUTED)
        add_text(s, x + 0.1, y + 0.1, widths[c_idx] - 0.2, 0.4, cell, size=10, bold=is_total, color=INK)
        x += widths[c_idx]
add_footer(s, 5)


# ---------------------------------------------------------------- Slide 6: Gap analysis
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "Where the gaps are", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "Your current cover vs the benchmarks for someone in your life stage.", size=14, color=MUTED)

headers = ["Coverage", "Benchmark", "Current", "Gap"]
widths = [3.5, 3.0, 2.5, 3.2]
x = 0.5
for i, h in enumerate(headers):
    add_box(s, x, 2.1, widths[i], 0.5, fill=INK)
    add_text(s, x + 0.1, 2.15, widths[i] - 0.2, 0.4, h, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    x += widths[i]

gaps = [
    ["Death / TPD", "10x income = $1.8M", "[$300K]", "−$1.5M", True],
    ["Major CI", "5x income + $100K = $1M", "[$0]", "−$1M", True],
    ["Hospital", "Private + rider", "[Standard, no rider]", "$5K+ co-pay/claim", True],
    ["Disability income", "60–70% of income", "[$0]", "Full income exposure", True],
    ["Wealth accumulation", "Structured plan", "[$80K savings]", "No engine", False],
]
for r_idx, row in enumerate(gaps):
    y = 2.6 + r_idx * 0.7
    is_critical = row[4]
    x = 0.5
    for c_idx, cell in enumerate(row[:4]):
        fill = LIGHT if r_idx % 2 == 0 else WHITE
        if c_idx == 3:
            fill = RGBColor(0xFD, 0xE2, 0xE2) if is_critical else RGBColor(0xFF, 0xF4, 0xD6)
        add_box(s, x, y, widths[c_idx], 0.7, fill=fill, line=MUTED)
        color = RED if c_idx == 3 and is_critical else INK
        add_text(s, x + 0.1, y + 0.18, widths[c_idx] - 0.2, 0.5, cell, size=11, bold=(c_idx == 3), color=color)
        x += widths[c_idx]
add_footer(s, 6)


# ---------------------------------------------------------------- Slide 7: Cost of doing nothing
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "Doing nothing isn't a neutral choice", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "What the current path costs over time.", size=14, color=MUTED)

# Big number callout
add_box(s, 0.7, 2.3, 5.8, 4.0, fill=RED)
add_text(s, 0.7, 2.6, 5.8, 0.6, "Cumulative cost of inaction", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.7, 3.3, 5.8, 1.4, "[$327K]", size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.7, 4.7, 5.8, 0.5, "to age 85", size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.7, 5.5, 5.8, 0.6, "[Premiums escalate to $1,000+/month\nin retirement years]", size=12, color=WHITE, align=PP_ALIGN.CENTER)

# Right side: explanation
add_text(s, 7.0, 2.5, 5.7, 0.6, "What this number is", size=18, bold=True, color=INK)
add_text(s, 7.0, 3.2, 5.7, 3.0,
    "[Cumulative hospital plan payments to 85 with 5% medical inflation]\n\n"
    "[Premiums today: $X/month]\n"
    "[Premiums at 75: $Y/month]\n"
    "[Premiums at 85: $Z/month]\n\n"
    "[Total over the period: $327K — paid TO the insurer, not back to you]",
    size=13, color=MUTED)
add_footer(s, 7)


# ---------------------------------------------------------------- Slide 8: Optional chart slide
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "The trajectory of inaction", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "Cumulative hospital plan premium outflow if nothing changes.", size=14, color=MUTED)

# Placeholder chart frame
add_box(s, 1.5, 2.2, 10.3, 4.6, fill=LIGHT)
add_text(s, 1.5, 4.3, 10.3, 0.5, "[Insert chart: cumulative cost vs age, anchored at age 85 = $327K]",
         size=16, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 4.85, 10.3, 0.5, "Use the tradingeconomics.com source data or AIA's hospital-plan inflation projection.",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)
add_footer(s, 8)


# ---------------------------------------------------------------- Slide 9: Option A
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "Option A — [Convert 100% of the expense]", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "What this option does, in one paragraph.", size=14, color=MUTED)

add_box(s, 0.7, 2.1, 6.0, 4.6, fill=LIGHT)
add_text(s, 1.0, 2.3, 5.5, 0.5, "What changes", size=16, bold=True, color=RED)
benefits = [
    "[1.] Generates income to offset the hospital plan premiums",
    "[2.] Provides an additional layer of passive income",
    "[3.] Capital can be preserved and passed to family on death",
    "[4.] Cash value accessible in later years if needed",
]
for i, b in enumerate(benefits):
    add_text(s, 1.0, 2.9 + i * 0.65, 5.7, 0.6, b, size=12, color=INK)

add_box(s, 7.0, 2.1, 5.7, 4.6, fill=ACCENT)
add_text(s, 7.3, 2.3, 5.2, 0.5, "Concept illustration", size=16, bold=True, color=INK)
add_text(s, 7.3, 2.9, 5.2, 0.5, "Capital required:  [$X]", size=14, color=INK)
add_text(s, 7.3, 3.5, 5.2, 0.5, "Income at 65:  [$Y/month]", size=14, color=INK)
add_text(s, 7.3, 4.1, 5.2, 0.5, "Bequest at 85:  [$Z]", size=14, color=INK)
add_text(s, 7.3, 4.8, 5.2, 1.7, "[One sentence explaining what this looks like in practice.]",
         size=12, color=MUTED)
add_footer(s, 9)


# ---------------------------------------------------------------- Slide 10: Option B
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "Option B — [Convert 2/3 of the expense]", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "Same shape as Option A, smaller capital commitment, more flexibility.", size=14, color=MUTED)

add_box(s, 0.7, 2.1, 6.0, 4.6, fill=LIGHT)
add_text(s, 1.0, 2.3, 5.5, 0.5, "What changes", size=16, bold=True, color=RED)
benefits = [
    "[1.] Same income generation logic, smaller scale",
    "[2.] Frees up [$X] capital for other uses",
    "[3.] Can be increased later (life-stage benefit if applicable)",
    "[4.] Lower commitment, easier to start now",
]
for i, b in enumerate(benefits):
    add_text(s, 1.0, 2.9 + i * 0.65, 5.7, 0.6, b, size=12, color=INK)

add_box(s, 7.0, 2.1, 5.7, 4.6, fill=ACCENT)
add_text(s, 7.3, 2.3, 5.2, 0.5, "Concept illustration", size=16, bold=True, color=INK)
add_text(s, 7.3, 2.9, 5.2, 0.5, "Capital required:  [0.67 × $X]", size=14, color=INK)
add_text(s, 7.3, 3.5, 5.2, 0.5, "Income at 65:  [0.67 × $Y/month]", size=14, color=INK)
add_text(s, 7.3, 4.1, 5.2, 0.5, "Bequest at 85:  [0.67 × $Z]", size=14, color=INK)
add_text(s, 7.3, 4.8, 5.2, 1.7, "[One sentence on flexibility vs Option A.]",
         size=12, color=MUTED)
add_footer(s, 10)


# ---------------------------------------------------------------- Slide 11: Side-by-side
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "Side-by-side — your call", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "Same outcome shape, different commitment level. There's no wrong answer.", size=14, color=MUTED)

headers = ["Dimension", "Option A (100%)", "Option B (2/3)"]
widths = [4.5, 4.0, 4.0]
x = 0.5
for i, h in enumerate(headers):
    add_box(s, x, 2.1, widths[i], 0.55, fill=INK)
    add_text(s, x + 0.1, 2.18, widths[i] - 0.2, 0.4, h, size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    x += widths[i]

rows = [
    ["Capital required", "[$327K]", "[$218K]"],
    ["Monthly income at 65", "[$X/month]", "[0.67X/month]"],
    ["Bequest at 85", "[$Y]", "[0.67Y]"],
    ["Cash flexibility", "Lower", "Higher"],
    ["Sleep-at-night factor", "[Higher]", "[Highest]"],
]
for r_idx, row in enumerate(rows):
    y = 2.65 + r_idx * 0.7
    x = 0.5
    for c_idx, cell in enumerate(row):
        fill = LIGHT if r_idx % 2 == 0 else WHITE
        add_box(s, x, y, widths[c_idx], 0.7, fill=fill, line=MUTED)
        bold = c_idx == 0
        add_text(s, x + 0.15, y + 0.18, widths[c_idx] - 0.3, 0.5, cell, size=12, bold=bold, color=INK)
        x += widths[c_idx]
add_footer(s, 11)


# ---------------------------------------------------------------- Slide 12: Bequest vs Dividends concept
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "Bequest vs Dividends — preserved AND paying", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4,
         "You can withdraw dividends for as long as you live, while preserving the bequest value at minimum the principal you deposited.",
         size=14, color=MUTED)

# Two columns
add_box(s, 0.7, 2.4, 6.0, 4.0, fill=LIGHT)
add_text(s, 1.0, 2.65, 5.5, 0.5, "Withdrawing dividends", size=18, bold=True, color=RED)
add_text(s, 1.0, 3.3, 5.5, 3.0,
         "Annualised dividend yield ~7.09% since inception\n\n"
         "Distributions are paid as long as the underlying fund is in force\n\n"
         "Dividend payout history is the substantiation — see slide 13",
         size=12, color=INK)

add_box(s, 7.0, 2.4, 5.7, 4.0, fill=ACCENT)
add_text(s, 7.3, 2.65, 5.2, 0.5, "Preserving the bequest", size=18, bold=True, color=INK)
add_text(s, 7.3, 3.3, 5.2, 3.0,
         "On death or terminal illness:\n"
         "  · Insurer pays the higher of\n"
         "    (a) 100% of total premiums paid, or\n"
         "    (b) the account value\n\n"
         "Capital protection is built in while the policy is in force.",
         size=12, color=INK)
add_footer(s, 12)


# ---------------------------------------------------------------- Slide 13: Track record
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "Track record — substantiation", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "This isn't an untested concept. The fund has paid dividends since inception.", size=14, color=MUTED)

# Big stat
add_box(s, 0.7, 2.2, 4.0, 4.5, fill=RED)
add_text(s, 0.7, 2.7, 4.0, 0.6, "Annualised dividend yield", size=12, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.7, 3.3, 4.0, 1.5, "~7.09%", size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.7, 5.0, 4.0, 0.5, "since inception", size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.7, 5.7, 4.0, 0.6, "[Source: fund factsheet]", size=10, color=WHITE, align=PP_ALIGN.CENTER)

# Right side
add_text(s, 5.2, 2.4, 7.5, 0.5, "What this means in practice", size=18, bold=True, color=INK)
add_text(s, 5.2, 3.0, 7.5, 4.0,
         "[The dividend payout has been resilient through:]\n"
         "  · [Multiple interest-rate cycles]\n"
         "  · [Equity market drawdowns (2020, 2022)]\n"
         "  · [Currency volatility across underlying holdings]\n\n"
         "[The track record is what gives confidence that this isn't a back-tested\n"
         "projection — it's an observed payout history.]\n\n"
         "[Past performance does not guarantee future outcomes — see disclaimer.]",
         size=12, color=INK)
add_footer(s, 13)


# ---------------------------------------------------------------- Slide 14: Underlying instrument
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "What you'd actually be investing in", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "Value proposition + underlying funds + positioning. One slide max.", size=14, color=MUTED)

# Three cards
items = [
    ("Value proposition", "[Income generation with bequest preservation\nthrough a multi-asset dividend mandate]"),
    ("Underlying funds", "[Allianz Income and Growth Fund (~70% of mandate)\n+ complementary income-tilted holdings]"),
    ("Fund positioning", "[Bond-heavy with equity dividend overlay;\ninverse relationship to interest-rate cycle]"),
]
card_w = 3.95
for i, (h, content) in enumerate(items):
    x = 0.6 + i * (card_w + 0.15)
    add_box(s, x, 2.2, card_w, 4.6, fill=LIGHT)
    add_box(s, x, 2.2, card_w, 0.6, fill=RED)
    add_text(s, x, 2.28, card_w, 0.5, h, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 3.0, card_w - 0.4, 3.5, content, size=12, color=INK)
add_footer(s, 14)


# ---------------------------------------------------------------- Slide 15: Timing rationale
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "Why now — the timing rationale", size=28, bold=True, color=INK)
add_accent_strip(s, 0.5, 1.05, 1.5, RED)
add_text(s, 0.5, 1.4, 12, 0.4, "Optional slide — only include if there's a real timing reason.", size=14, color=MUTED)

add_box(s, 0.7, 2.1, 12.0, 4.6, fill=LIGHT)
add_text(s, 1.2, 2.4, 11.0, 0.5, "Interest-rate cycle", size=16, bold=True, color=RED)
add_text(s, 1.2, 3.0, 11.0, 3.5,
         "[For bond-heavy income funds, fund price has an INVERSE relationship to interest rates:]\n\n"
         "  · [When rates rise → fund price drops (entry timing favourable)]\n"
         "  · [When rates fall → fund price rises (entry timing less favourable)]\n\n"
         "[Singapore interest rates are currently on the higher end of the 20-year range. With rate cuts incoming,\n"
         "this is a favourable entry window for a dividend-income mandate.]\n\n"
         "[Source: tradingeconomics.com/singapore/interest-rate]",
         size=13, color=INK)
add_footer(s, 15)


# ---------------------------------------------------------------- Slide 16: Disclaimer
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.4, 12, 0.6, "Disclaimer", size=24, bold=True, color=INK)
add_accent_strip(s, 0.5, 0.95, 1.5, RED)

disclaimer = (
    "All diagrams and charts included in this proposal serve illustrative purposes only. Actual results may differ "
    "based on the performance of the underlying assets in the selected funds. Fund performance is not guaranteed, "
    "and the policy's cash value may fall below the total premiums invested.\n\n"
    "For precise projections across various premium levels, please refer to the official policy illustration. All "
    "figures shown are rounded to the nearest dollar and do not reflect potential changes in net asset value (NAV). "
    "Investment management fees and product charges are estimated on a best-effort basis. Actual payable amounts may vary.\n\n"
    "Capital protection applies in the event of death or terminal illness. If the policy is in force, the insurer will "
    "pay the higher of 100% of total premiums paid or the account value.\n\n"
    "Past performance does not represent or predict future outcomes, and asset values may fluctuate with market conditions. "
    "Early withdrawals may result in charges. Please speak with a licensed financial advisor representative before making "
    "any financial decisions.\n\n"
    "Any observations, recommendations, or projections in this document reflect opinions only and should not be "
    "interpreted as financial advice. For the most accurate and up-to-date details, always consult the official product "
    "materials, websites, and contract documents.\n\n"
    "This document is private and confidential. It is intended solely for the recipient and may not be copied, distributed, "
    "reproduced, or shared with any third party, in whole or in part."
)
add_text(s, 0.5, 1.3, 12.3, 5.7, disclaimer, size=10, color=MUTED)
add_footer(s, 16)


# ---------------------------------------------------------------- Slide 17: Thank you + Next steps
s = prs.slides.add_slide(BLANK)
add_text(s, 0.5, 0.5, 12, 0.8, "Thank you.", size=48, bold=True, color=RED)
add_text(s, 0.5, 1.5, 12, 0.5, "Next steps", size=22, bold=True, color=INK)

steps = [
    "Relationship comes first — no decision today is the right answer if it's not yours",
    "[Specific next action — e.g. review existing policy summaries within 7 days]",
    "[Time-anchored milestone — e.g. submit application by [date] for the Vitality discount lock-in]",
    "Yearly review built into the relationship",
]
for i, st in enumerate(steps):
    y = 2.1 + i * 0.6
    add_text(s, 0.6, y, 0.4, 0.5, "·", size=22, bold=True, color=RED)
    add_text(s, 1.0, y, 11.5, 0.5, st, size=15, color=INK)

# Contact card
add_box(s, 0.5, 5.4, 12.3, 1.5, fill=LIGHT)
add_text(s, 0.8, 5.55, 11.5, 0.5, "[Your name] · AIA Singapore", size=18, bold=True, color=INK)
add_text(s, 0.8, 6.05, 11.5, 0.4, "LinkedIn: [linkedin.com/in/...]   ·   Mobile: [+65 ...]   ·   Email: [you@aia.com.sg]",
         size=13, color=MUTED)
add_footer(s, 17)


# Save
prs.save(str(OUT))
print(f"wrote {OUT} ({OUT.stat().st_size / 1024:.1f} KB, 17 slides)")
